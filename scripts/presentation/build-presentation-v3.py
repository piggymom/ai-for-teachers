#!/usr/bin/env python3
"""
Build AI Literacy Presentation — clean from-scratch build.
7 slides, NV dark navy branding, Calibri, speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── BRAND ──
DARK = RGBColor(0x1A, 0x1F, 0x36)
CARD = RGBColor(0x24, 0x2A, 0x45)
INDIGO = RGBColor(0x36, 0x45, 0x9C)
TEAL = RGBColor(0x0E, 0x96, 0x8A)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE = RGBColor(0xE8, 0x6C, 0x3A)
PERI = RGBColor(0xCA, 0xCF, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0x8A, 0x8F, 0xA8)
BGRAY = RGBColor(0x5A, 0x5F, 0x78)
FNT = "Calibri"


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK


def tb(slide, l, t, w, h, text, sz=16, c=WHITE, b=False, a=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.name = FNT
    p.font.bold = b
    p.alignment = a
    return tf


def ap(tf, text, sz=16, c=WHITE, b=False, a=PP_ALIGN.LEFT, sp=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.name = FNT
    p.font.bold = b
    p.alignment = a
    p.space_before = Pt(sp)
    return p


def rr(slide, l, t, w, h, fill, text="", sz=14, fc=WHITE, b=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        s.text_frame.paragraphs[0].text = text
        s.text_frame.paragraphs[0].font.size = Pt(sz)
        s.text_frame.paragraphs[0].font.color.rgb = fc
        s.text_frame.paragraphs[0].font.name = FNT
        s.text_frame.paragraphs[0].font.bold = b
    return s


def bar(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def dot(slide, x, y, sz, fill, text="", fsz=11):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        s.text_frame.paragraphs[0].text = text
        s.text_frame.paragraphs[0].font.size = Pt(fsz)
        s.text_frame.paragraphs[0].font.color.rgb = WHITE
        s.text_frame.paragraphs[0].font.bold = True
        s.text_frame.paragraphs[0].font.name = FNT
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ── CREATE ──
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
blank = prs.slide_layouts[6]


# ════════════════════════════════════════════
# SLIDE 1 — THE META OPEN
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
bar(s, Inches(0), Inches(0), Inches(10), Inches(0.06), TEAL)
tb(s, Inches(0.8), Inches(1.3), Inches(8.4), Inches(1),
   "Critical AI Literacy", 44, WHITE, True)
tb(s, Inches(0.8), Inches(2.4), Inches(8.4), Inches(0.5),
   "A practice, not a credential.", 20, PERI)
bar(s, Inches(0.8), Inches(3.3), Inches(2.2), Inches(0.04), TEAL)
tb(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.5),
   "New Visions for Public Schools", 15, LGRAY)

notes(s, """THE META OPEN (~3 min)
Minimal slide. This is you, not the deck.

"I want to start with something unusual for a work PD. I'm going to tell you where I actually am with AI — not where I think I'm supposed to be.

Over the past couple of years I've built things with AI. Real things that real educators use. A feedback tool with over a hundred weekly users. A PD course with an AI tutor built in. And in building those things I've made mistakes — not just technical ones, but ethical ones. Moments where I shipped something and then thought: wait, should I have done that?

That experience sent me back to the reading. The OECD's new AI literacy draft from May, the Center for Digital Thriving framework from Harvard, and five days ago the Department of Labor released its first-ever AI Literacy Framework — February 13th.

AI literacy isn't a credential you earn once. It's a practice. I'm in it. You're in it. And today I want us to do some of that work together rather than me just presenting at you."
""")


# ════════════════════════════════════════════
# SLIDE 2 — WHAT IS AI ACTUALLY DOING
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3),
   "HOW IT WORKS", 12, TEAL, True)
tb(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6),
   "What Is AI Actually Doing?", 32, WHITE, True)

# Flow diagram
labels = ["Your\nprompt", "Tokens", "Weights +\nProbabilities", "Next likely\ntoken", "Output"]
colors = [INDIGO, VIOLET, ORANGE, TEAL, INDIGO]
bw, bh = Inches(1.5), Inches(0.85)
sx, y = Inches(0.5), Inches(2.0)
gap = Inches(0.38)

for i, (lab, clr) in enumerate(zip(labels, colors)):
    x = sx + i * (bw + gap)
    rr(s, x, y, bw, bh, clr, lab, 12, WHITE, True)
    if i < 4:
        tb(s, x + bw + Inches(0.05), y + Inches(0.2), Inches(0.28), Inches(0.4),
           "\u2192", 20, LGRAY, a=PP_ALIGN.CENTER)

# Insight box
rr(s, Inches(0.5), Inches(3.4), Inches(9.0), Inches(1.7), CARD)
tf = tb(s, Inches(0.8), Inches(3.55), Inches(8.4), Inches(1.4),
        "Pattern-matching at scale that produces outputs that look like understanding.",
        15, PERI)
ap(tf, "Confidence is a design choice, not an indicator of accuracy.", 13, LGRAY, sp=10)
ap(tf, "Bias is baked in \u2014 the training data came from us.", 13, LGRAY, sp=4)

notes(s, """WHAT IS AI ACTUALLY DOING (~3 min)
OECD K1.1-K1.3 in 60 seconds.

"When you type into ChatGPT or Claude, your words get broken into tokens. Those tokens get processed through billions of numerical weights set during training on human text.

At every step, it asks: what token is most likely next? It's a sophisticated prediction machine. Not reasoning. Not understanding. Pattern-matching at scale.

When AI sounds confident, that's a design choice. It generates authoritative text whether correct or not. Knowing that changes how you read the output.

And those weights came from training on human text — human biases, human assumptions, and the labor of often underpaid workers who labeled data are baked into every response."
""")


# ════════════════════════════════════════════
# SLIDE 3 — THE SPEED OF CHANGE
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3),
   "CONTEXT", 12, TEAL, True)
tb(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6),
   "The Speed of Change", 32, WHITE, True)

eras = [
    ("1950s\u201380s", "Rule-Based AI", "Humans write\nevery rule"),
    ("1990s\u20132000s", "Machine\nLearning", "Systems learn\nfrom data"),
    ("2010s", "Deep\nLearning", "Neural nets,\nimage & speech"),
    ("2020\u201322", "Large Language\nModels", "Scale +\ntransformers"),
    ("Now", "Agentic AI", "Models plan,\nact, use tools"),
]
nc = [BGRAY, INDIGO, VIOLET, ORANGE, TEAL]
cw, ch = Inches(1.65), Inches(2.2)
cx = Inches(0.45)
cy = Inches(1.6)
cg = Inches(0.2)

for i, (era, name, desc) in enumerate(eras):
    x = cx + i * (cw + cg)
    rr(s, x, cy, cw, ch, CARD)
    bar(s, x, cy, cw, Inches(0.06), nc[i])
    tb(s, x + Inches(0.1), cy + Inches(0.2), cw - Inches(0.2), Inches(0.3),
       era, 10, LGRAY, a=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.1), cy + Inches(0.55), cw - Inches(0.2), Inches(0.7),
       name, 14, WHITE, True, PP_ALIGN.CENTER)
    tb(s, x + Inches(0.1), cy + Inches(1.35), cw - Inches(0.2), Inches(0.7),
       desc, 11, LGRAY, a=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, x + cw + Inches(0.02), cy + Inches(0.85), Inches(0.16), Inches(0.4),
           "\u203A", 20, BGRAY, a=PP_ALIGN.CENTER)

tb(s, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.8),
   "Our ethics, policies, and instincts developed for a much slower landscape. "
   "We are genuinely making it up as we go.", 14, PERI)

notes(s, """THE SPEED OF CHANGE (~1 min)
Quick, don't linger. CUT THIS SLIDE if running short — it's the least load-bearing.

"For decades, AI meant rule-based systems. Then machine learning. Deep learning. And then — in a handful of years — LLMs and now agentic AI that can plan and take actions.

Our ethics frameworks developed for a much slower landscape. We are genuinely making it up as we go. Which means rooms like this one actually matter."
""")


# ════════════════════════════════════════════
# SLIDE 4 — THE GAP
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3),
   "THE GAP", 12, ORANGE, True)
tb(s, Inches(0.8), Inches(0.55), Inches(8.4), Inches(0.55),
   "What Most AI Literacy Frameworks Miss", 28, WHITE, True)

rows = [
    ("Technical Literacy", "How does AI work? How do I use it?",
     "Covered by most frameworks", INDIGO, TEAL),
    ("Market Literacy", "Who built this? What are their incentives?",
     "Largely missing", ORANGE, ORANGE),
    ("Self-Literacy", "Do I know my values well enough to notice the shift?",
     "Almost entirely absent", VIOLET, VIOLET),
]

gy = Inches(1.4)
rh = Inches(1.05)

for i, (label, desc, status, bc, tc) in enumerate(rows):
    y = gy + i * (rh + Inches(0.15))
    rr(s, Inches(0.5), y, Inches(9.0), rh, CARD)
    bar(s, Inches(0.5), y, Inches(0.08), rh, bc)
    tb(s, Inches(0.8), y + Inches(0.2), Inches(2.2), Inches(0.6),
       label, 15, WHITE, True)
    tb(s, Inches(3.1), y + Inches(0.25), Inches(3.3), Inches(0.5),
       desc, 12, LGRAY)
    tb(s, Inches(6.6), y + Inches(0.25), Inches(2.7), Inches(0.5),
       status, 11, tc, True, PP_ALIGN.RIGHT)

tb(s, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.5),
   "DOL Framework (Feb 2026): Understand \u2022 Explore \u2022 Direct \u2022 Evaluate \u2022 Use Responsibly  \u2014  maps almost entirely to Row 1",
   11, LGRAY)

notes(s, """THE GAP IN AI LITERACY (~3 min)

"Most AI literacy focuses on technical literacy. CDT at Harvard adds two more layers:

Market literacy — who built this, what are their incentives, what data is being collected about you.

Self-literacy — do you know your own values clearly enough to notice when AI is quietly nudging you away from them?

The DOL just released a national framework — five content areas that cover technical literacy well. But essentially nothing about market or self-literacy.

That's not a criticism — it's a workforce framework doing workforce things. But it shows the gap.

The question isn't whether you can use AI. It's whether you can use AI and still be recognizably yourself at the end of it."
""")


# ════════════════════════════════════════════
# SLIDE 5 — ACTIVITY SETUP
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3),
   "INTERACTIVE ACTIVITY", 12, TEAL, True)
tb(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6),
   "Align on the Line", 32, WHITE, True)

# Scenario card
rr(s, Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.3), CARD)
tf = tb(s, Inches(1.1), Inches(1.85), Inches(7.8), Inches(1.0),
        "You're applying for your dream job. It's competitive.", 17, WHITE)
ap(tf, "The application asks for a cover letter, work samples, "
   "and a short video introduction.", 17, WHITE, sp=8)

# Scale
tb(s, Inches(1.2), Inches(3.45), Inches(2.0), Inches(0.3),
   "Totally Fine", 14, TEAL, True)
tb(s, Inches(6.8), Inches(3.45), Inches(2.0), Inches(0.3),
   "Crosses a Line", 14, ORANGE, True, PP_ALIGN.RIGHT)
bar(s, Inches(1.2), Inches(3.85), Inches(7.6), Inches(0.08), BGRAY)

dc = [TEAL, INDIGO, VIOLET, ORANGE, ORANGE]
for i in range(5):
    dx = Inches(1.2) + i * Inches(1.9)
    dot(s, dx - Inches(0.06), Inches(3.78), Inches(0.2), dc[i])

tb(s, Inches(0.8), Inches(4.3), Inches(8.4), Inches(0.8),
   "Drag each AI use to where it sits on your personal scale.\n"
   "Individual first. No judgment. Go with your gut.",
   14, LGRAY, a=PP_ALIGN.CENTER)

notes(s, """ACTIVITY SETUP (~1 min of 10-min activity)

"You're applying for your dream job. It's competitive. The application asks for a cover letter, work samples, and a short video introduction.

On your Miro board you'll see eight different ways someone might use AI in that process. Drag each one to where it sits on your personal scale. Totally Fine on the left, Crosses a Line on the right.

Individual first. No judgment. Go with your gut. You've got about two and a half minutes."
""")


# ════════════════════════════════════════════
# SLIDE 6 — THE 8 AI USES
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3),
   "ALIGN ON THE LINE", 12, TEAL, True)
tb(s, Inches(0.8), Inches(0.55), Inches(8.4), Inches(0.45),
   "The 8 AI Uses", 24, WHITE, True)

uses = [
    "Using AI to check grammar and spelling",
    "Asking AI to suggest stronger word choices",
    "Giving AI your draft and asking it to improve the flow",
    "Giving AI the job description and asking it to rewrite your letter to match",
    "Asking AI to write the full letter from scratch based on your resume",
    "Using AI to generate talking points for your video",
    "Using an AI voice/video tool to make your delivery sound more polished",
    "Asking AI to research the hiring manager and tailor your letter to their preferences",
]
uc = [TEAL, TEAL, INDIGO, INDIGO, VIOLET, VIOLET, ORANGE, ORANGE]
colw = Inches(4.1)
cardh = Inches(0.72)
xo = [Inches(0.5), Inches(5.0)]
ys = Inches(1.2)

for i, use in enumerate(uses):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = xo[col]
    y = ys + row * (cardh + Inches(0.18))

    rr(s, x, y, colw, cardh, CARD)
    dot(s, x + Inches(0.12), y + (cardh - Inches(0.3)) // 2,
        Inches(0.3), uc[i], str(i + 1), 11)
    tb(s, x + Inches(0.5), y + Inches(0.08),
       colw - Inches(0.65), cardh - Inches(0.16), use, 11, WHITE)

tb(s, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.5),
   "What is the underlying principle driving YOUR line?  "
   "Authenticity \u2022 Fairness \u2022 Accountability",
   13, PERI, a=PP_ALIGN.CENTER)

notes(s, """THE 8 AI USES — DEBRIEF (~8 min)

"Let's look at this together. Look at number four — 'give AI the job description and ask it to rewrite your letter to match.' We've got spread on that one. Who wants to say what put them where they are?"

(Take 2-3 responses. Don't resolve. Let the tension sit.)

"And number eight — researching the hiring manager. What's the line there? Privacy? Fairness? What you'd want done to you?"

At ~8 min:
"What is the underlying principle driving your line? Authenticity? Fairness? Accountability? None are wrong. But they lead to very different lines.

And now — imagine your students doing exactly this. For their college essays. Right now. Tonight. That's why we're having this conversation."
""")


# ════════════════════════════════════════════
# SLIDE 7 — FROM PERSONAL TO INSTITUTIONAL
# ════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
tb(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3),
   "CLOSING", 12, VIOLET, True)
tb(s, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.55),
   "From Personal to Institutional", 32, WHITE, True)

qs = [
    ("Where does AI amplify our work\nwithout eroding what makes it ours?", TEAL),
    ("What are we being asked to be\ntransparent about \u2014 and to whom?", ORANGE),
    ("Who isn't in this conversation\nthat should be?", VIOLET),
]

for i, (q, clr) in enumerate(qs):
    y = Inches(1.6) + i * Inches(1.15)
    dot(s, Inches(0.8), y + Inches(0.05), Inches(0.45), clr, str(i + 1), 16)
    tb(s, Inches(1.5), y, Inches(7.5), Inches(0.8), q, 18, WHITE)

# Bottom provocation
rr(s, Inches(0.5), Inches(4.55), Inches(9.0), Inches(0.75), CARD)
tf = tb(s, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.6),
        "Every tool we adopt was built by a company with investors and a growth model.",
        12, PERI, a=PP_ALIGN.CENTER)
ap(tf, "What are we trading \u2014 and did we decide that was a fair trade?",
   12, PERI, a=PP_ALIGN.CENTER, sp=4)

notes(s, """FROM PERSONAL TO INSTITUTIONAL (~3 min)

"I don't want to give you a policy today. I want to leave you with three questions:

1. Where does AI amplify our work without eroding what makes it ours?
2. What are we being asked to be transparent about, and to whom?
3. Who isn't in this conversation that should be?

The DOL framework calls on education systems to integrate AI literacy. The question is whether we adopt the framework as given, or decide our communities deserve the fuller version — including market and self-literacy.

Every tool we're encouraged to adopt was built by a company with investors and a growth model. What are we trading? And did we decide that was a fair trade?"

CLOSE:
"I'll drop a reading list in the chat. The OECD draft, the CDT framework (one page, reframes everything), and the DOL AI Literacy Framework. Thanks for being in this with me."

READING LIST:
- OECD AI Literacy Framework (May 2025)
- CDT Agency Framework (Harvard, one-page grid)
- U.S. DOL AI Literacy Framework TEN 07-25 (Feb 13, 2026)
""")


# ── SAVE ──
out = "/Users/carladelporto/Downloads/AI_Literacy_Critical_Presentation.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
