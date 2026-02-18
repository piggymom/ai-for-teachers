#!/usr/bin/env python3
"""Build AI Literacy Presentation using existing NV deck as template base."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
DARK_NAVY = RGBColor(0x1A, 0x1F, 0x36)
CARD_NAVY = RGBColor(0x24, 0x2A, 0x45)
INDIGO = RGBColor(0x36, 0x45, 0x9C)
TEAL = RGBColor(0x0E, 0x96, 0x8A)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE = RGBColor(0xE8, 0x6C, 0x3A)
LIGHT_PERIWINKLE = RGBColor(0xCA, 0xCF, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x8A, 0x8F, 0xA8)
BLUE_GRAY = RGBColor(0x5A, 0x5F, 0x78)
FONT = "Calibri"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def clear_placeholders(slide):
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def tbox(slide, left, top, width, height, text, size=16, color=WHITE,
         bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.bold = bold
    p.alignment = align
    return tf


def add_p(tf, text, size=16, color=WHITE, bold=False, align=PP_ALIGN.LEFT, before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.bold = bold
    p.alignment = align
    p.space_before = Pt(before)
    return p


def rect(slide, l, t, w, h, fill, text="", size=14, fc=WHITE, bold=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        s.text_frame.word_wrap = True
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        s.text_frame.paragraphs[0].text = text
        s.text_frame.paragraphs[0].font.size = Pt(size)
        s.text_frame.paragraphs[0].font.color.rgb = fc
        s.text_frame.paragraphs[0].font.name = FONT
        s.text_frame.paragraphs[0].font.bold = bold
    return s


def circle(slide, x, y, sz, fill, text="", size=11):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if text:
        s.text_frame.paragraphs[0].text = text
        s.text_frame.paragraphs[0].font.size = Pt(size)
        s.text_frame.paragraphs[0].font.color.rgb = WHITE
        s.text_frame.paragraphs[0].font.bold = True
        s.text_frame.paragraphs[0].font.name = FONT
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# Load existing NV deck as template, delete its slides via lxml
from lxml import etree
template = "/Users/carladelporto/Library/CloudStorage/GoogleDrive-asher.rhys.scott@gmail.com/My Drive/AI_All_Staff_Feb2026_Redesigned.pptx"
prs = Presentation(template)

# Delete all existing slides properly
for _ in range(len(prs.slides)):
    sldId = prs.slides._sldIdLst[0]
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        # Try without namespace
        for attr in sldId.attrib:
            if attr.endswith('}id') or attr == 'r:id':
                rId = sldId.attrib[attr]
                break
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(sldId)

blank = prs.slide_layouts[6]

# ── SLIDE 1: THE META OPEN ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
rect(s, Inches(0), Inches(0), Inches(10), Inches(0.06), TEAL)
tbox(s, Inches(0.8), Inches(1.3), Inches(8.4), Inches(1), "Critical AI Literacy", 44, WHITE, True)
tbox(s, Inches(0.8), Inches(2.4), Inches(8.4), Inches(0.5), "A practice, not a credential.", 20, LIGHT_PERIWINKLE)
rect(s, Inches(0.8), Inches(3.3), Inches(2.2), Inches(0.04), TEAL)
tbox(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.5), "New Visions for Public Schools", 15, LIGHT_GRAY)
notes(s, "SLIDE 1 \u2014 THE META OPEN (~3 min)\n\n\"I want to start with something unusual for a work PD. I'm going to tell you where I actually am with AI \u2014 not where I think I'm supposed to be.\n\nOver the past couple of years I've built things with AI. Real things that real educators use. A feedback tool with over a hundred weekly users. A PD course with an AI tutor built in. And in building those things I've made mistakes \u2014 not just technical ones, but ethical ones.\n\nThat experience sent me back to the reading. The OECD's new AI literacy draft, the Center for Digital Thriving framework from Harvard, and five days ago the Department of Labor released its first-ever AI Literacy Framework.\n\nAI literacy isn't a credential you earn once. It's a practice. I'm in it. You're in it. Today I want us to do some of that work together.\"")

# ── SLIDE 2: WHAT IS AI ACTUALLY DOING ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3), "HOW IT WORKS", 12, TEAL, True)
tbox(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6), "What Is AI Actually Doing?", 32, WHITE, True)

labels = ["Your prompt", "Tokens", "Weights +\nProbabilities", "Next likely\ntoken", "Output"]
colors = [INDIGO, VIOLET, ORANGE, TEAL, INDIGO]
bw, bh = Inches(1.5), Inches(0.85)
sx, y = Inches(0.5), Inches(2.0)
gap = Inches(0.38)
for i, (lab, clr) in enumerate(zip(labels, colors)):
    x = sx + i * (bw + gap)
    rect(s, x, y, bw, bh, clr, lab, 13, WHITE, True)
    if i < 4:
        tbox(s, x + bw + Inches(0.05), y + Inches(0.2), Inches(0.28), Inches(0.4), "\u2192", 20, LIGHT_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(0.5), Inches(3.4), Inches(9.0), Inches(1.7), CARD_NAVY)
tf = tbox(s, Inches(0.8), Inches(3.55), Inches(8.4), Inches(1.4),
          "Pattern-matching at scale that produces outputs that look like understanding.", 15, LIGHT_PERIWINKLE)
add_p(tf, "Confidence is a design choice, not an indicator of accuracy. Bias is baked in \u2014 the training data came from us.", 13, LIGHT_GRAY, before=10)
notes(s, "SLIDE 2 \u2014 WHAT IS AI ACTUALLY DOING (~3 min)\n\nYour words \u2192 tokens \u2192 processed through billions of weights \u2192 predicts next token \u2192 output.\n\nIt is not reasoning. It is pattern-matching at scale. Confidence is a design choice. Bias from human training data is baked in.")

# ── SLIDE 3: EVOLUTION TIMELINE ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3), "CONTEXT", 12, TEAL, True)
tbox(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6), "The Speed of Change", 32, WHITE, True)

eras = [
    ("1950s\u201380s", "Rule-Based AI", "Humans write\nevery rule"),
    ("1990s\u20132000s", "Machine\nLearning", "Systems learn\nfrom data"),
    ("2010s", "Deep\nLearning", "Neural nets,\nimage & speech"),
    ("2020\u201322", "Large Language\nModels", "Scale +\ntransformers"),
    ("Now", "Agentic AI", "Models plan,\nact, use tools"),
]
nc = [BLUE_GRAY, INDIGO, VIOLET, ORANGE, TEAL]
cw, ch = Inches(1.65), Inches(2.2)
cx, cy = Inches(0.45), Inches(1.6)
cg = Inches(0.2)

for i, (era, name, desc) in enumerate(eras):
    x = cx + i * (cw + cg)
    rect(s, x, cy, cw, ch, CARD_NAVY)
    rect(s, x, cy, cw, Inches(0.06), nc[i])
    tbox(s, x + Inches(0.1), cy + Inches(0.2), cw - Inches(0.2), Inches(0.3), era, 10, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    tbox(s, x + Inches(0.1), cy + Inches(0.55), cw - Inches(0.2), Inches(0.7), name, 14, WHITE, True, PP_ALIGN.CENTER)
    tbox(s, x + Inches(0.1), cy + Inches(1.35), cw - Inches(0.2), Inches(0.7), desc, 11, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 4:
        tbox(s, x + cw + Inches(0.02), cy + Inches(0.85), Inches(0.16), Inches(0.4), "\u203A", 20, BLUE_GRAY, align=PP_ALIGN.CENTER)

tbox(s, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.8),
     "Our ethics, policies, and instincts developed for a much slower landscape. We are genuinely making it up as we go.", 14, LIGHT_PERIWINKLE)
notes(s, "SLIDE 3 \u2014 TIMELINE (~1 min). Cut this slide if short on time. It's the least load-bearing.")

# ── SLIDE 4: THE GAP ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3), "THE GAP", 12, ORANGE, True)
tbox(s, Inches(0.8), Inches(0.55), Inches(8.4), Inches(0.55), "What Most AI Literacy Frameworks Miss", 28, WHITE, True)

rows = [
    ("Technical Literacy", "How does AI work? How do I use it?", "Covered by most frameworks", INDIGO, TEAL),
    ("Market Literacy", "Who built this? What are their incentives?", "Largely missing", ORANGE, ORANGE),
    ("Self-Literacy", "Do I know my values well enough to notice the shift?", "Almost entirely absent", VIOLET, VIOLET),
]
gy, rh, rg = Inches(1.4), Inches(1.05), Inches(0.15)
for i, (lab, desc, tag, bc, tc) in enumerate(rows):
    y = gy + i * (rh + rg)
    rect(s, Inches(0.5), y, Inches(9.0), rh, CARD_NAVY)
    rect(s, Inches(0.5), y, Inches(0.08), rh, bc)
    tbox(s, Inches(0.8), y + Inches(0.2), Inches(2.2), Inches(0.6), lab, 15, WHITE, True)
    tbox(s, Inches(3.1), y + Inches(0.25), Inches(3.3), Inches(0.5), desc, 12, LIGHT_GRAY)
    tbox(s, Inches(6.6), y + Inches(0.25), Inches(2.7), Inches(0.5), tag, 11, tc, True, PP_ALIGN.RIGHT)

tbox(s, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.5),
     "DOL Framework (Feb 2026): Understand \u2022 Explore \u2022 Direct \u2022 Evaluate \u2022 Use Responsibly  \u2014  maps almost entirely to Row 1", 11, LIGHT_GRAY)
notes(s, "SLIDE 4 \u2014 THE GAP (~3 min)\n\nCDT adds Market Literacy and Self-Literacy. DOL's 5 content areas cover technical literacy well but miss the other two dimensions.")

# ── SLIDE 5: ALIGN ON THE LINE SETUP ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.3), "INTERACTIVE ACTIVITY", 12, TEAL, True)
tbox(s, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.6), "Align on the Line", 32, WHITE, True)
rect(s, Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.3), CARD_NAVY)
tf = tbox(s, Inches(1.1), Inches(1.85), Inches(7.8), Inches(1.0),
          "You're applying for your dream job. It's competitive.", 17, WHITE)
add_p(tf, "The application asks for a cover letter, work samples, and a short video introduction.", 17, WHITE, before=8)

tbox(s, Inches(1.2), Inches(3.5), Inches(2.0), Inches(0.3), "Totally Fine", 14, TEAL, True)
tbox(s, Inches(6.8), Inches(3.5), Inches(2.0), Inches(0.3), "Crosses a Line", 14, ORANGE, True, PP_ALIGN.RIGHT)
rect(s, Inches(1.2), Inches(3.85), Inches(7.6), Inches(0.12), BLUE_GRAY)
for i in range(5):
    dx = Inches(1.2) + i * Inches(1.9)
    circle(s, dx - Inches(0.08), Inches(3.78), Inches(0.25), [TEAL, INDIGO, VIOLET, ORANGE, ORANGE][i])

tbox(s, Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.8),
     "Drag each AI use to where it sits on your personal scale.\nIndividual first. No judgment. Go with your gut.", 14, LIGHT_GRAY, align=PP_ALIGN.CENTER)
notes(s, "SLIDE 5 \u2014 ACTIVITY SETUP. Give 2.5 minutes for individual placement on Miro board.")

# ── SLIDE 6: THE 8 USES ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3), "ALIGN ON THE LINE", 12, TEAL, True)
tbox(s, Inches(0.8), Inches(0.55), Inches(8.4), Inches(0.45), "The 8 AI Uses", 24, WHITE, True)

uses = [
    "Using AI to check grammar and spelling",
    "Asking AI to suggest stronger word choices",
    "Giving AI your draft and asking it to improve the flow",
    "Giving AI the job description and asking it to rewrite your letter to match",
    "Asking AI to write the full letter from scratch based on your resume",
    "Using AI to generate talking points for your video",
    "Using an AI voice/video tool to make your delivery sound more polished",
    "Asking AI to research the hiring manager and tailor your letter to their preferences",
]
uc = [TEAL, TEAL, INDIGO, INDIGO, VIOLET, VIOLET, ORANGE, ORANGE]
cw2, cardh = Inches(4.1), Inches(0.72)
xo = [Inches(0.5), Inches(5.0)]
ys, yg2 = Inches(1.2), Inches(0.18)

for i, use in enumerate(uses):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = xo[col]
    y = ys + row * (cardh + yg2)
    rect(s, x, y, cw2, cardh, CARD_NAVY)
    circle(s, x + Inches(0.12), y + (cardh - Inches(0.3)) / 2, Inches(0.3), uc[i], str(i+1), 11)
    tbox(s, x + Inches(0.5), y + Inches(0.08), cw2 - Inches(0.65), cardh - Inches(0.16), use, 11, WHITE)

tbox(s, Inches(0.5), Inches(4.85), Inches(9.0), Inches(0.5),
     "What is the underlying principle driving YOUR line?  Authenticity \u2022 Fairness \u2022 Accountability",
     13, LIGHT_PERIWINKLE, align=PP_ALIGN.CENTER)
notes(s, "SLIDE 6 \u2014 DEBRIEF (~8 min)\n\nCall out #4 and #8 for discussion. Ask: what principle drives your line? Then pivot: imagine your students doing this for college essays. Tonight.")

# ── SLIDE 7: FROM PERSONAL TO INSTITUTIONAL ──
s = prs.slides.add_slide(blank)
set_bg(s, DARK_NAVY)
clear_placeholders(s)
tbox(s, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3), "CLOSING", 12, VIOLET, True)
tbox(s, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.55), "From Personal to Institutional", 32, WHITE, True)

qs = [
    ("Where does AI amplify our work\nwithout eroding what makes it ours?", TEAL),
    ("What are we being asked to be\ntransparent about \u2014 and to whom?", ORANGE),
    ("Who isn't in this conversation\nthat should be?", VIOLET),
]
for i, (q, clr) in enumerate(qs):
    y = Inches(1.6) + i * Inches(1.15)
    circle(s, Inches(0.8), y + Inches(0.05), Inches(0.45), clr, str(i+1), 16)
    tbox(s, Inches(1.5), y, Inches(7.5), Inches(0.8), q, 18, WHITE)

rect(s, Inches(0.5), Inches(4.55), Inches(9.0), Inches(0.75), CARD_NAVY)
tbox(s, Inches(0.8), Inches(4.6), Inches(8.4), Inches(0.6),
     "Every tool we adopt was built by a company with investors and a growth model.\nWhat are we trading \u2014 and did we decide that was a fair trade?",
     12, LIGHT_PERIWINKLE, align=PP_ALIGN.CENTER)
notes(s, "SLIDE 7 \u2014 CLOSING (~3 min)\n\nThree questions for New Visions. End with reading list: OECD draft, CDT framework, DOL AI Literacy Framework (TEN 07-25).")

# Save
out = "/Users/carladelporto/Downloads/AI_Literacy_Critical_Presentation.pptx"
prs.save(out)
print(f"Saved: {out} ({len(prs.slides)} slides)")
