#!/usr/bin/env python3
"""
Build the AI Literacy Presentation for New Visions.
7 slides, dark navy NV branding, Calibri font, speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# =============================================================================
# BRAND CONSTANTS
# =============================================================================

# Backgrounds
DARK_NAVY = RGBColor(0x1A, 0x1F, 0x36)
CARD_NAVY = RGBColor(0x24, 0x2A, 0x45)
LIGHT_BG = RGBColor(0xF5, 0xF6, 0xFA)

# Accents
INDIGO = RGBColor(0x36, 0x45, 0x9C)
TEAL = RGBColor(0x0E, 0x96, 0x8A)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE = RGBColor(0xE8, 0x6C, 0x3A)
LIGHT_PERIWINKLE = RGBColor(0xCA, 0xCF, 0xF0)

# Text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x8A, 0x8F, 0xA8)
BLUE_GRAY = RGBColor(0x5A, 0x5F, 0x78)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)

FONT = "Calibri"

SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


# =============================================================================
# HELPERS
# =============================================================================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_paragraph(tf, text, font_size=16, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=6, space_after=4,
                  font_name=FONT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # no border
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.name = FONT
        tf.paragraphs[0].font.bold = bold
        shape.text_frame.margin_left = Inches(0.15)
        shape.text_frame.margin_right = Inches(0.15)
        shape.text_frame.margin_top = Inches(0.1)
        shape.text_frame.margin_bottom = Inches(0.1)
    return shape


def add_line(slide, start_x, start_y, end_x, end_y, color=INDIGO, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def set_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# =============================================================================
# BUILD PRESENTATION
# =============================================================================

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]

# ─────────────────────────────────────────────
# SLIDE 1 — THE META OPEN
# ─────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, DARK_NAVY)

# Thin accent bar at top
add_rounded_rect(slide1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), TEAL)

# Title
add_text_box(slide1, Inches(0.8), Inches(1.3), Inches(8.4), Inches(1.2),
             "Critical AI Literacy", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.LEFT)

# Subtitle
tf1 = add_text_box(slide1, Inches(0.8), Inches(2.4), Inches(8.4), Inches(0.6),
                   "A practice, not a credential.", font_size=20, color=LIGHT_PERIWINKLE,
                   alignment=PP_ALIGN.LEFT)

# Author / org
add_text_box(slide1, Inches(0.8), Inches(3.6), Inches(8.4), Inches(0.8),
             "New Visions for Public Schools", font_size=15, color=LIGHT_GRAY,
             alignment=PP_ALIGN.LEFT)

# Thin accent line
add_line(slide1, Inches(0.8), Inches(3.45), Inches(3.0), Inches(3.45), TEAL, Pt(2))

set_notes(slide1, """SLIDE 1 — THE META OPEN (~3 minutes)
Visual: Minimal. Your face or a blank slide. This is you, not a deck.
Facilitator note: The meta-layer is the point — you're modeling what AI literacy actually looks like as a practice, not a credential.

SCRIPT:
"I want to start with something a little unusual for a work PD. I'm going to tell you where I actually am with AI — not where I think I'm supposed to be.

Over the past couple of years I've built things with AI. Real things that real educators use. A feedback tool that now has over a hundred weekly users. A professional development course with an AI tutor built into it. And in building those things I've made a lot of mistakes — not just technical ones, but ethical ones. Moments where I shipped something and then thought: wait, should I have done that? Who does this serve? What am I actually asking people to trust?

That experience sent me back to the reading. Right now I'm working through the OECD's new AI literacy draft — which came out in May — and a framework from the Center for Digital Thriving at Harvard that I think is genuinely the most honest framing of AI literacy I've seen. And five days ago the Department of Labor released its first-ever AI Literacy Framework — a federal definition of what AI literacy means for the American workforce. It just came out February 13th. I'll show you how it maps against the other frameworks in a minute.

I'm sharing all of this because I think the most important thing I can model today is that AI literacy isn't a credential you earn once. It's a practice. I'm in it. You're in it. And today I want us to do some of that work together rather than me just presenting at you."
""")


# ─────────────────────────────────────────────
# SLIDE 2 — WHAT IS AI ACTUALLY DOING
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, DARK_NAVY)

# Section label
add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.4),
             "HOW IT WORKS", font_size=12, color=TEAL, bold=True)

# Title
add_text_box(slide2, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.8),
             "What Is AI Actually Doing?", font_size=32, color=WHITE, bold=True)

# Flow diagram — 5 boxes with arrows
labels = ["Your prompt", "Tokens", "Weights +\nProbabilities", "Next likely\ntoken", "Output"]
colors = [INDIGO, VIOLET, ORANGE, TEAL, INDIGO]
box_w = Inches(1.55)
box_h = Inches(0.85)
start_x = Inches(0.5)
y = Inches(2.1)
gap = Inches(0.35)

for i, (label, clr) in enumerate(zip(labels, colors)):
    x = start_x + i * (box_w + gap)
    add_rounded_rect(slide2, x, y, box_w, box_h, clr, label,
                     font_size=13, font_color=WHITE, bold=True)
    # Arrow between boxes
    if i < len(labels) - 1:
        arrow_x = x + box_w + Inches(0.02)
        arrow_y = y + box_h / 2
        arrow_end = arrow_x + gap - Inches(0.04)
        add_line(slide2, arrow_x, arrow_y, arrow_end, arrow_y, LIGHT_GRAY, Pt(2))

# Key insight box at bottom
add_rounded_rect(slide2, Inches(0.5), Inches(3.5), Inches(9.0), Inches(1.6), CARD_NAVY)
tf2 = add_text_box(slide2, Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.4),
                   "Pattern-matching at scale that produces outputs that look like understanding.",
                   font_size=15, color=LIGHT_PERIWINKLE, bold=False)
add_paragraph(tf2, "Confidence is a design choice, not an indicator of accuracy. "
              "Bias is baked in — the training data came from us.",
              font_size=13, color=LIGHT_GRAY, space_before=8)

set_notes(slide2, """SLIDE 2 — WHAT IS AI ACTUALLY DOING (~3 minutes)
Design note: OECD K1.1–K1.3 knowledge in 60 seconds with one diagram.

SCRIPT:
"Before we can have an honest ethical conversation about AI, we need a shared mental model of what it's actually doing. So here's the 60-second version.

When you type something into ChatGPT or Claude, your words get broken into fragments called tokens — think of them as chunks of text, roughly a word or part of a word. Those tokens get processed through billions of numerical weights — values that were set during training by exposing the model to an enormous amount of human-generated text.

What the model is doing, at every step, is asking: given everything before this point, what token is most likely to come next? It's a very sophisticated prediction machine. It is not reasoning. It is not understanding. It is pattern-matching at a scale that produces outputs that look like understanding.

Here's why that matters practically: when AI sounds confident, that's a design choice. The model doesn't have an 'I'm not sure about this' instinct built in by default. It generates fluent, authoritative-sounding text whether it's correct or not. Knowing that changes how you read the output.

And here's the part people skip: those weights were set by training on human text — which means human biases, human assumptions, and the labor of often underpaid human workers who labeled data are baked into every response you get. Numbers in, numbers out — but the numbers came from us."
""")


# ─────────────────────────────────────────────
# SLIDE 3 — THE EVOLUTION TIMELINE
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, DARK_NAVY)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.4),
             "CONTEXT", font_size=12, color=TEAL, bold=True)

add_text_box(slide3, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.8),
             "The Speed of Change", font_size=32, color=WHITE, bold=True)

# Timeline: horizontal line with 5 nodes
eras = [
    ("1950s–80s", "Rule-Based AI", "Humans write\nevery rule"),
    ("1990s–2000s", "Machine\nLearning", "Systems learn\nfrom data"),
    ("2010s", "Deep\nLearning", "Neural networks,\nimage & speech"),
    ("2020–22", "Large Language\nModels", "Scale +\ntransformers"),
    ("Now", "Agentic AI", "Models that plan,\nact, use tools"),
]

line_y = Inches(2.8)
start_x = Inches(0.6)
end_x = Inches(9.4)
add_line(slide3, start_x, line_y, end_x, line_y, BLUE_GRAY, Pt(2))

node_colors = [BLUE_GRAY, INDIGO, VIOLET, ORANGE, TEAL]

for i, (era, name, desc) in enumerate(eras):
    x = start_x + i * ((end_x - start_x) / (len(eras) - 1))
    # Node dot
    dot_size = Inches(0.2)
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL, x - dot_size // 2, line_y - dot_size // 2,
        dot_size, dot_size
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = node_colors[i]
    dot.line.fill.background()

    # Era label above
    add_text_box(slide3, x - Inches(0.6), line_y - Inches(1.0), Inches(1.2), Inches(0.3),
                 era, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Name
    add_text_box(slide3, x - Inches(0.75), line_y - Inches(0.7), Inches(1.5), Inches(0.5),
                 name, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description below
    add_text_box(slide3, x - Inches(0.75), line_y + Inches(0.2), Inches(1.5), Inches(0.6),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_text_box(slide3, Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.8),
             "Our ethics, policies, and instincts developed for a much slower landscape.\nWe are genuinely making it up as we go.",
             font_size=14, color=LIGHT_PERIWINKLE, alignment=PP_ALIGN.LEFT)

set_notes(slide3, """SLIDE 3 — THE EVOLUTION TIMELINE (~1 minute)
Facilitator note: Quick, don't linger. If you run short on total time, cut this slide. It's the least load-bearing.

SCRIPT:
"Just to orient us historically — because I think the speed of this matters.

For decades, AI meant rule-based systems. Humans wrote every rule. Then machine learning meant systems could find patterns in data. Deep learning gave us image recognition, speech. And then — in basically a handful of years — we got large language models that can generate text, and now we're into agentic AI that can plan and take actions in the world.

The reason I show this isn't to be dazzling. It's to say: our ethics frameworks, our policies, our instincts — they developed for a much slower moving landscape. We are genuinely making it up as we go. Which means rooms like this one actually matter."
""")


# ─────────────────────────────────────────────
# SLIDE 4 — THE GAP IN MOST AI LITERACY WORK
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, DARK_NAVY)

add_text_box(slide4, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.4),
             "THE GAP", font_size=12, color=ORANGE, bold=True)

add_text_box(slide4, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.6),
             "What Most AI Literacy Frameworks Miss",
             font_size=28, color=WHITE, bold=True)

# CDT 3-row grid on the left
row_labels = ["Technical Literacy", "Market Literacy", "Self-Literacy"]
row_descs = [
    "How does AI work?\nHow do I use it?",
    "Who built this?\nWhat are their incentives?",
    "Do I know my values well\nenough to notice the shift?"
]
row_colors = [INDIGO, ORANGE, VIOLET]
coverage = ["Covered by most frameworks", "Largely missing", "Almost entirely absent"]

grid_x = Inches(0.5)
grid_y = Inches(1.4)
row_h = Inches(1.1)
label_w = Inches(2.2)
desc_w = Inches(3.5)
tag_w = Inches(2.8)

for i in range(3):
    y_pos = grid_y + i * (row_h + Inches(0.15))

    # Color bar on left
    slide4.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, grid_x, y_pos, Inches(0.08), row_h
    ).fill.solid()
    slide4.shapes[-1].fill.fore_color.rgb = row_colors[i]
    slide4.shapes[-1].line.fill.background()

    # Row background
    add_rounded_rect(slide4, grid_x + Inches(0.12), y_pos, Inches(8.9), row_h, CARD_NAVY)

    # Label
    add_text_box(slide4, grid_x + Inches(0.3), y_pos + Inches(0.15), label_w, Inches(0.7),
                 row_labels[i], font_size=15, color=WHITE, bold=True)

    # Description
    add_text_box(slide4, grid_x + Inches(2.5), y_pos + Inches(0.15), desc_w, Inches(0.7),
                 row_descs[i], font_size=12, color=LIGHT_GRAY)

    # Coverage tag
    tag_color = TEAL if i == 0 else ORANGE if i == 1 else VIOLET
    add_text_box(slide4, grid_x + Inches(6.2), y_pos + Inches(0.3), tag_w, Inches(0.4),
                 coverage[i], font_size=11, color=tag_color, bold=True,
                 alignment=PP_ALIGN.RIGHT)

# DOL sidebar
add_text_box(slide4, Inches(0.8), Inches(4.85), Inches(8.4), Inches(0.5),
             "DOL Framework (Feb 2026): Understand \u2022 Explore \u2022 Direct \u2022 Evaluate \u2022 Use Responsibly  \u2014  maps almost entirely to Row 1",
             font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

set_notes(slide4, """SLIDE 4 — THE GAP IN MOST AI LITERACY WORK (~3 minutes)
Visual: CDT 3-row grid with DOL mapping.

DOL's 5 Foundational Content Areas:
1. Understand AI Principles
2. Explore AI Uses
3. Direct AI Effectively
4. Evaluate AI Outputs
5. Use AI Responsibly

CDT's 3 literacy dimensions: Technical, Market, Self

SCRIPT:
"So there's a lot of AI literacy work out there right now — the OECD framework I mentioned, UNESCO, ISTE, all of it. And most of it focuses here — on technical literacy. How does AI work, how do you use it, how do you evaluate its outputs. That stuff is necessary. But the Center for Digital Thriving at Harvard argues it's not sufficient.

They add two more layers. Market literacy — understanding the business models behind these tools. Who built this, what are their incentives, what data is being collected about you — not for you. Because every free AI tool is a company with a revenue model, and it's worth knowing what you're trading before you build your whole workflow around it.

And then self-literacy — do you actually know your own values clearly enough to notice when AI is quietly nudging you away from them? That one's uncomfortable. It's the piece most PD skips entirely.

And to make this concrete — the Department of Labor just released a national AI Literacy Framework last week. Five content areas. And they're good: understand how AI works, explore uses, learn to prompt well, evaluate outputs, use it responsibly. If you map those onto CDT's grid, they cover the technical literacy row thoroughly. The 'Use AI Responsibly' area touches ethics — it talks about protecting sensitive data, following workplace policies, maintaining accountability. But there's essentially nothing about market literacy — who built these tools, what's the revenue model, what data is being collected about you. And nothing about self-literacy — whether you know your own values well enough to notice when they're being shifted.

That's not a criticism of DOL — it's a workforce framework doing workforce things. But it shows you the gap. The federal government is telling us AI literacy matters. CDT is telling us the version most people are teaching leaves out the parts that require the most judgment.

I think the question isn't whether you can use AI. It's whether you can use AI and still be recognizably yourself at the end of it. That's what we're going to explore right now."
""")


# ─────────────────────────────────────────────
# SLIDE 5 — ALIGN ON THE LINE: SETUP
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, DARK_NAVY)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.4),
             "INTERACTIVE ACTIVITY", font_size=12, color=TEAL, bold=True)

add_text_box(slide5, Inches(0.8), Inches(0.7), Inches(8.4), Inches(0.7),
             "Align on the Line", font_size=32, color=WHITE, bold=True)

# Scenario card
add_rounded_rect(slide5, Inches(0.8), Inches(1.7), Inches(8.4), Inches(1.4), CARD_NAVY)

scenario_tf = add_text_box(slide5, Inches(1.1), Inches(1.85), Inches(7.8), Inches(1.1),
                           "You're applying for your dream job. It's competitive.",
                           font_size=17, color=WHITE, bold=False)
add_paragraph(scenario_tf, "The application asks for a cover letter, work samples, "
              "and a short video introduction.",
              font_size=17, color=WHITE, space_before=8)

# Scale visual
scale_y = Inches(3.5)
add_line(slide5, Inches(1.5), scale_y, Inches(8.5), scale_y, LIGHT_GRAY, Pt(3))

# Left label
add_text_box(slide5, Inches(0.8), scale_y - Inches(0.15), Inches(1.5), Inches(0.3),
             "Totally Fine", font_size=12, color=TEAL, bold=True, alignment=PP_ALIGN.LEFT)

# Right label
add_text_box(slide5, Inches(7.7), scale_y - Inches(0.15), Inches(1.8), Inches(0.3),
             "Crosses a Line", font_size=12, color=ORANGE, bold=True, alignment=PP_ALIGN.RIGHT)

# Scale dots
for i in range(5):
    dx = Inches(1.5) + i * (Inches(7.0) / 4)
    dot = slide5.shapes.add_shape(MSO_SHAPE.OVAL, dx - Inches(0.06), scale_y - Inches(0.06),
                                  Inches(0.12), Inches(0.12))
    dot.fill.solid()
    c = [TEAL, INDIGO, VIOLET, ORANGE, ORANGE][i]
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

# Instruction
add_text_box(slide5, Inches(0.8), Inches(4.1), Inches(8.4), Inches(0.8),
             "Drag each AI use to where it sits on your personal scale.\nIndividual first. No judgment. Go with your gut.",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

set_notes(slide5, """SLIDE 5 — ALIGN ON THE LINE: SETUP (~1 minute of 10-minute activity)

SCRIPT:
"OK, here's the activity. I want you to imagine something real.

You're applying for your dream job. It's competitive — maybe it's a stretch role, maybe it's something you've wanted for years. The application asks for a cover letter, a couple of work samples, and a short video introduction.

On your Miro board you'll see eight different ways someone might use AI in that process — from pretty light touch to pretty significant. I want you to drag each one to where it sits on your personal scale. Totally Fine on the left, Crosses a Line on the right, and everything in between.

This is individual first. No judgment. Go with your gut. You've got about two and a half minutes."
""")


# ─────────────────────────────────────────────
# SLIDE 6 — ALIGN ON THE LINE: THE 8 USES
# ─────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, DARK_NAVY)

add_text_box(slide6, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.3),
             "ALIGN ON THE LINE", font_size=12, color=TEAL, bold=True)

add_text_box(slide6, Inches(0.8), Inches(0.55), Inches(8.4), Inches(0.5),
             "The 8 AI Uses", font_size=24, color=WHITE, bold=True)

uses = [
    "Using AI to check grammar and spelling",
    "Asking AI to suggest stronger word choices",
    "Giving AI your draft and asking it to improve the flow",
    "Giving AI the job description and asking it to rewrite your letter to match",
    "Asking AI to write the full letter from scratch based on your resume",
    "Using AI to generate talking points for your video",
    "Using an AI voice/video tool to make your delivery sound more polished",
    "Asking AI to research the hiring manager and tailor your letter to their preferences",
]

# 2 columns, 4 rows
col_w = Inches(4.1)
card_h = Inches(0.72)
x_offsets = [Inches(0.5), Inches(5.0)]
y_start = Inches(1.2)
y_gap = Inches(0.18)

# Color gradient from teal (light touch) to orange (heavy)
use_colors = [TEAL, TEAL, INDIGO, INDIGO, VIOLET, VIOLET, ORANGE, ORANGE]

for i, use in enumerate(uses):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = x_offsets[col]
    y = y_start + row * (card_h + y_gap)

    add_rounded_rect(slide6, x, y, col_w, card_h, CARD_NAVY)

    # Number badge
    badge_size = Inches(0.3)
    badge = slide6.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.12), y + (card_h - badge_size) / 2,
        badge_size, badge_size
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = use_colors[i]
    badge.line.fill.background()
    badge.text_frame.paragraphs[0].text = str(i + 1)
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = FONT
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Use text
    add_text_box(slide6, x + Inches(0.5), y + Inches(0.08), col_w - Inches(0.65), card_h - Inches(0.16),
                 use, font_size=11, color=WHITE)

# Bottom prompt
add_text_box(slide6, Inches(0.8), Inches(4.85), Inches(8.4), Inches(0.5),
             "What is the underlying principle driving YOUR line?  Authenticity \u2022 Fairness \u2022 Accountability",
             font_size=13, color=LIGHT_PERIWINKLE, alignment=PP_ALIGN.CENTER)

set_notes(slide6, """SLIDE 6 — ALIGN ON THE LINE: THE 8 USES (debrief portion, ~8 minutes)

Debrief script:
"OK, let's look at this together. I'm going to call out a few that look interesting...

Look at number four — 'give AI the job description and ask it to rewrite your letter to match it.' We've got people spread pretty much across the whole scale on that one. Who wants to say what put them where they are on that one?"

(Take 2–3 responses. Don't resolve it. Let the tension sit.)

"And look at number eight — researching the hiring manager. For some of you that crosses a line. I'm curious what the line is there. Is it about privacy? About fairness? About what you'd want done to you?"

(1–2 more responses.)

At around 8 minutes:
"Here's the question I actually want to pull out of all of this: what is the underlying principle that's driving your line?

For some of you it might be about authenticity — this should represent me. For others it might be about fairness — other applicants aren't using this, so I shouldn't either. For others it might be about accountability — I'm willing to use it as long as I can stand behind the result.

None of those are wrong. But notice that they lead to very different lines. And now —

(pause for effect)

— imagine your students doing exactly this. For their college essays. Right now. Tonight. That's why we're having this conversation."
""")


# ─────────────────────────────────────────────
# SLIDE 7 — FROM PERSONAL TO INSTITUTIONAL
# ─────────────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, DARK_NAVY)

add_text_box(slide7, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.4),
             "CLOSING", font_size=12, color=VIOLET, bold=True)

add_text_box(slide7, Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.6),
             "From Personal to Institutional", font_size=32, color=WHITE, bold=True)

questions = [
    ("Where does AI amplify our work\nwithout eroding what makes it ours?", TEAL),
    ("What are we being asked to be\ntransparent about — and to whom?", ORANGE),
    ("Who isn't in this conversation\nthat should be?", VIOLET),
]

q_y = Inches(1.6)
for i, (q, clr) in enumerate(questions):
    y_pos = q_y + i * Inches(1.15)

    # Number
    num_size = Inches(0.45)
    num_shape = slide7.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), y_pos + Inches(0.05), num_size, num_size
    )
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = clr
    num_shape.line.fill.background()
    num_shape.text_frame.paragraphs[0].text = str(i + 1)
    num_shape.text_frame.paragraphs[0].font.size = Pt(16)
    num_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].font.name = FONT
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Question text
    add_text_box(slide7, Inches(1.5), y_pos, Inches(7.5), Inches(0.8),
                 q, font_size=18, color=WHITE, bold=False)

# Bottom provocation
add_rounded_rect(slide7, Inches(0.5), Inches(4.6), Inches(9.0), Inches(0.7), CARD_NAVY)
add_text_box(slide7, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.5),
             "Every tool we adopt was built by a company with investors and a growth model.\nWhat are we trading — and did we decide that was a fair trade?",
             font_size=12, color=LIGHT_PERIWINKLE, alignment=PP_ALIGN.CENTER)

set_notes(slide7, """SLIDE 7 — FROM PERSONAL TO INSTITUTIONAL (~3 minutes)

SCRIPT:
"I don't want to give you a policy today. I think that's the wrong move at this stage. What I want to leave you with is three questions I think New Visions genuinely needs to sit with.

First: where does AI amplify our work without eroding what makes it ours? There are real answers to that question — things we do that are repetitive, administrative, time-consuming, where AI genuinely frees us up to do the human work better. I've seen it. But 'it can help' isn't the same as 'we've thought carefully about where.'

Second: what are we being asked to be transparent about, and to whom? Teachers using AI to write feedback. Administrators using it to draft communications. There are people on the receiving end of those things who may want to know. Where do we have an obligation to say so?

And third — and this is the CDT collective agency question: who isn't in this conversation that should be? We're all staff. But the people most affected by AI decisions in schools are students and families. What does it mean that we're setting norms they haven't been part of?

The DOL framework explicitly calls on education systems to integrate AI literacy — and it names K-12, training providers, and state agencies as audiences. So this isn't theoretical. Federal guidance is now pointed at organizations like ours. The question is whether we adopt the framework as given, or whether we decide our students and communities deserve the fuller version — the one that includes the market and self-literacy dimensions too.

I'll end with the market literacy provocation that I think about a lot: every tool we're being encouraged to adopt was built by a company with investors and a growth model. Before we build institutional workflows around it, it's worth asking — what are we trading? And did we decide that was a fair trade, or did it just kind of happen?"

CLOSE:
"I'll drop a short reading list in the chat. The OECD draft is genuinely worth reading. The CDT framework is one page and it reframes everything. The DOL AI Literacy Framework is brand new and worth knowing about — it's the federal government's first attempt to define this for the workforce. Thanks for being in this with me."

READING LIST:
- OECD AI Literacy Framework — Draft, May 2025
- Center for Digital Thriving Agency Framework — Harvard, one-page grid
- U.S. Department of Labor AI Literacy Framework — TEN 07-25, February 13, 2026
""")


# =============================================================================
# SAVE
# =============================================================================

output_path = "/Users/carladelporto/Library/CloudStorage/GoogleDrive-asher.rhys.scott@gmail.com/My Drive/AI_Literacy_Critical_Presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
