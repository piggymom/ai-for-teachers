#!/usr/bin/env python3
"""Build the AI for Teachers master pitch deck as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LIGHT = RGBColor(0x99, 0xF6, 0xE4)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
SLATE_LIGHT = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
ROSE = RGBColor(0xF4, 0x3F, 0x5E)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
VERY_LIGHT_TEAL = RGBColor(0xF0, 0xFD, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=SLATE, bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_para(tf, text, font_size=18, color=SLATE, bold=False, italic=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


# ─── SLIDE 1: HOOK ────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, SLATE)

# Accent bar at top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2)).text_frame
set_text(tf, "What happens when the AI goes away?", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tf2 = add_textbox(slide, Inches(2), Inches(3.5), Inches(9), Inches(1.5)).text_frame
set_text(tf2, "Students using AI tools showed a 17% performance drop\nwhen access was removed.", font_size=24, color=TEAL_LIGHT, alignment=PP_ALIGN.CENTER)

tf3 = add_textbox(slide, Inches(2), Inches(5.3), Inches(9), Inches(1)).text_frame
set_text(tf3, "Most AI tools for teachers create the same dependency.\nWe built something different.", font_size=20, color=SLATE_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "This is the 'sit up' moment. Don't explain yet — let the question land. Pause after reading the stat. Let the audience feel the tension before moving on.")


# ─── SLIDE 2: THE PROBLEM ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1)).text_frame
set_text(tf, "Teachers are drowning in AI anxiety, not AI tools", font_size=36, color=SLATE, bold=True)

# Three stat boxes
stats = [
    ("82%", "of K-12 teachers feel\nunprepared to use AI", ROSE),
    ("$50-200", "per teacher per PD session\n— and they call it 'useless'", AMBER),
    ("Feb 2026", "DOL AI Literacy Framework\n— districts must now comply", TEAL),
]

for i, (num, desc, accent) in enumerate(stats):
    left = Inches(0.8 + i * 4.1)
    # Accent line
    add_rect(slide, left, Inches(2.2), Inches(3.5), Inches(0.06), accent)
    # Number
    tf = add_textbox(slide, left, Inches(2.5), Inches(3.5), Inches(1.2)).text_frame
    set_text(tf, num, font_size=52, color=accent, bold=True)
    # Description
    tf2 = add_textbox(slide, left, Inches(3.9), Inches(3.5), Inches(1.5)).text_frame
    set_text(tf2, desc, font_size=18, color=SLATE_LIGHT)

# Bottom line
tf = add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(1)).text_frame
set_text(tf, "The gap: no structured, personalized, evidence-based program exists to build\nteacher AI capacity at scale.", font_size=20, color=SLATE, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Establish the triple bind: teachers are anxious, PD is broken, and now there's a federal mandate. Pause on each stat. The DOL framework is new — most audiences haven't heard of it yet.")


# ─── SLIDE 3: THE OZEMPIC INSIGHT ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLATE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Capacity vs. Dependency", font_size=36, color=WHITE, bold=True)

# Two columns
# Left: Dependency
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0x2D, 0x1B, 0x1B))
tf = add_textbox(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.6)).text_frame
set_text(tf, "THE DEPENDENCY MODEL", font_size=16, color=ROSE, bold=True)

dep_items = [
    "AI generates content FOR the teacher",
    "Teacher reviews and delivers",
    "Skill required: clicking buttons",
    "When tool is removed: capability disappears",
]
tf = add_textbox(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3)).text_frame
set_text(tf, dep_items[0], font_size=17, color=RGBColor(0xCC, 0xCC, 0xCC))
for item in dep_items[1:]:
    add_para(tf, item, font_size=17, color=RGBColor(0xCC, 0xCC, 0xCC), space_before=Pt(12))

# Right: Capacity
add_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0x0C, 0x2E, 0x2B))
tf = add_textbox(slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(0.6)).text_frame
set_text(tf, "THE CAPACITY MODEL", font_size=16, color=EMERALD, bold=True)

cap_items = [
    "Teacher LEARNS the skill with AI coaching",
    "Teacher practices in any AI tool",
    "Skill required: prompting, evaluation, iteration",
    "When tool is removed: teacher keeps the skill",
]
tf = add_textbox(slide, Inches(7.4), Inches(2.7), Inches(4.8), Inches(3)).text_frame
set_text(tf, cap_items[0], font_size=17, color=TEAL_LIGHT)
for item in cap_items[1:]:
    add_para(tf, item, font_size=17, color=TEAL_LIGHT, space_before=Pt(12))

# Bottom quote
tf = add_textbox(slide, Inches(1.5), Inches(6.5), Inches(10), Inches(0.8)).text_frame
set_text(tf, '"Most AI tools suppress the struggle that creates mastery. We build capacity, not dependency."', font_size=18, color=TEAL_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "This is the intellectual core of the pitch. Take your time here. The Ozempic analogy: the drug works while you're on it, but when you stop, the weight comes back — because the underlying behaviors never changed. Same with AI tools that do the work for teachers.")


# ─── SLIDE 4: THE SOLUTION ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Meet Skippy", font_size=42, color=TEAL_DARK, bold=True)

tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8)).text_frame
set_text(tf, "An AI tutor that teaches teachers to use AI — by building, not watching.", font_size=24, color=SLATE)

features = [
    ("7-week personalized course", "Structured progression from \"What is AI?\" to personal AI policy"),
    ("20 minutes per week", "Self-paced, asynchronous — no coverage needed, no scheduling"),
    ("One artifact every session", "Prompt templates, lesson workflows, feedback rubrics, differentiation strategies"),
    ("Skills transfer to any AI tool", "Teachers practice in ChatGPT and Gemini — not locked into Skippy"),
]

for i, (title, desc) in enumerate(features):
    top = Inches(2.6 + i * 1.1)
    add_rect(slide, Inches(0.8), top, Inches(0.08), Inches(0.7), TEAL)
    tf = add_textbox(slide, Inches(1.2), top, Inches(10), Inches(0.4)).text_frame
    set_text(tf, title, font_size=20, color=SLATE, bold=True)
    tf2 = add_textbox(slide, Inches(1.2), Emu(top + Inches(0.4)), Inches(10), Inches(0.4)).text_frame
    set_text(tf2, desc, font_size=16, color=SLATE_LIGHT)

add_notes(slide, "Keep this slide fast. The 'how' comes next. Emphasize 'by building, not watching' — that's the key phrase. The four features map to the four objections districts have: too long, too much time, no evidence, tool lock-in.")


# ─── SLIDE 5: THE 7-WEEK JOURNEY ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "One Win Per Week", font_size=36, color=SLATE, bold=True)

weeks = [
    ("0", "Onboarding", "Teaching Profile", "5 min"),
    ("1", "Understanding AI", "AI Mental Model Card", "~20 min"),
    ("2", "Prompting", "4C Prompt Template", "~20 min"),
    ("3", "Lesson Planning", "Planning Workflow", "~20 min"),
    ("4", "Feedback", "Feedback Template", "~20 min"),
    ("5", "Differentiation", "Differentiation Template", "~20 min"),
    ("6", "Ethics & Policy", "Personal AI Policy", "~20 min"),
]

# Header row
headers = ["WEEK", "FOCUS", "TEACHER BUILDS", "TIME"]
header_widths = [Inches(1), Inches(3), Inches(4.5), Inches(1.5)]
header_lefts = [Inches(1)]
for w in header_widths[:-1]:
    header_lefts.append(Emu(header_lefts[-1] + w))

top = Inches(1.5)
add_rect(slide, Inches(0.8), top, Inches(11.5), Inches(0.5), TEAL)
for j, h in enumerate(headers):
    tf = add_textbox(slide, header_lefts[j], top, header_widths[j], Inches(0.5)).text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    set_text(tf, h, font_size=14, color=WHITE, bold=True)

for i, (wk, focus, artifact, time) in enumerate(weeks):
    row_top = Emu(Inches(2.05) + i * Inches(0.65))
    bg_color = VERY_LIGHT_TEAL if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.8), row_top, Inches(11.5), Inches(0.6), bg_color)

    vals = [f"Week {wk}", focus, artifact, time]
    for j, val in enumerate(vals):
        tf = add_textbox(slide, header_lefts[j], row_top, header_widths[j], Inches(0.6)).text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        c = TEAL_DARK if j == 0 else SLATE
        b = j == 0
        set_text(tf, val, font_size=15, color=c, bold=b)

tf = add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.7)).text_frame
set_text(tf, "Every session produces something the teacher uses Monday morning. Not theory. Not slides. A tool they built.", font_size=17, color=SLATE_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Emphasize 'artifact-first' — every session produces something the teacher uses Monday morning. The portfolio of 7 artifacts is the proof of learning and the basis for PD credit documentation.")


# ─── SLIDE 6: THE LEARNING SCIENCE ────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "This Isn't Prompt Tips. It's Pedagogy.", font_size=36, color=SLATE, bold=True)

sciences = [
    ("SOLO Taxonomy", "Real-time diagnostic classifier assesses teacher understanding across 5 validated levels", TEAL),
    ("4C Framework", "Context, Constraints, Command, Criteria — the transferable prompting structure", AMBER),
    ("External Testing Loop", '"Try it in ChatGPT. Come back and tell me what happened."', EMERALD),
    ("One Win Then Wrap", "Every conversation follows a 6-phase arc toward one concrete output", RGBColor(0x63, 0x66, 0xF1)),
    ("Adaptive Scaffolding", "Skippy's behavior visibly changes based on diagnosed level — from analogies to peer-mode", ROSE),
]

for i, (name, desc, color) in enumerate(sciences):
    top = Inches(1.6 + i * 1.05)
    add_rect(slide, Inches(0.8), top, Inches(0.4), Inches(0.4), color)
    tf = add_textbox(slide, Inches(1.5), top, Inches(4), Inches(0.5)).text_frame
    set_text(tf, name, font_size=19, color=SLATE, bold=True)
    tf2 = add_textbox(slide, Inches(5.5), top, Inches(7), Inches(0.7)).text_frame
    set_text(tf2, desc, font_size=16, color=SLATE_LIGHT)

add_notes(slide, "This is where you win technical audiences. For foundations: spend time here — the SOLO taxonomy has 40+ years of validation. For districts: go fast, emphasize 'adapts to every teacher automatically.'")


# ─── SLIDE 7: PERSONALIZATION ENGINE ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLATE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "6 Layers of Personalization", font_size=36, color=WHITE, bold=True)

layers = [
    ("1", "Global Personality", "Warm British sensibility, expert tutor, artifact-first design"),
    ("2", "Week-Specific Pedagogy", "Learning goals, phase management, skill targets"),
    ("3", "Personalized Opening", "Name, subject, grade level woven into first message"),
    ("4", "Teacher Profile Context", "Goals, constraints, time drains, AI experience"),
    ("5", "Conversation Ledger", "Phase, SOLO level, 4C progress, artifact state"),
    ("6", "Diagnostic Probe", "Opening assessment calibrates every session"),
]

for i, (num, name, desc) in enumerate(layers):
    top = Inches(1.5 + i * 0.9)
    # Layer number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), top, Inches(0.55), Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    tf_c = shape.text_frame
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_text(tf_c, num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    tf = add_textbox(slide, Inches(2.1), top, Inches(3.5), Inches(0.5)).text_frame
    set_text(tf, name, font_size=18, color=WHITE, bold=True)
    tf2 = add_textbox(slide, Inches(6), top, Inches(6.5), Inches(0.5)).text_frame
    set_text(tf2, desc, font_size=15, color=TEAL_LIGHT)

tf = add_textbox(slide, Inches(1.5), Inches(6.8), Inches(10), Inches(0.6)).text_frame
set_text(tf, '"Every response is shaped by who you are, what you know, and where you are in the conversation."', font_size=17, color=TEAL_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "This slide proves depth. For tech audiences, this is the 'wow, they actually built something' moment. For non-technical audiences: 'Imagine a tutor who remembers everything about your teaching context and adjusts every response.'")


# ─── SLIDE 8: DEMO / SCREENSHOTS ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Show, Don't Tell", font_size=36, color=SLATE, bold=True)

# Placeholder boxes for screenshots
screens = [
    ("Personalized Onboarding", "5-min profile captures subject,\ngrade, goals, and constraints"),
    ("Adaptive Conversation", "Skippy responds differently\nbased on SOLO level"),
    ("Artifact Extraction", "Teacher's work saved\nautomatically to gallery"),
    ("Dashboard & Progress", "Portfolio of artifacts +\nweek-by-week completion"),
]

for i, (title, desc) in enumerate(screens):
    left = Inches(0.6 + i * 3.15)
    # Screenshot placeholder
    shape = add_rect(slide, left, Inches(1.8), Inches(2.9), Inches(3.5), WHITE)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)

    tf_inner = add_textbox(slide, left, Inches(2.8), Inches(2.9), Inches(1.5)).text_frame
    set_text(tf_inner, "[Screenshot]", font_size=24, color=RGBColor(0xCB, 0xD5, 0xE1), alignment=PP_ALIGN.CENTER)

    tf = add_textbox(slide, left, Inches(5.5), Inches(2.9), Inches(0.5)).text_frame
    set_text(tf, title, font_size=15, color=TEAL_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    tf2 = add_textbox(slide, left, Inches(6.0), Inches(2.9), Inches(0.8)).text_frame
    set_text(tf2, desc, font_size=13, color=SLATE_LIGHT, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Replace placeholder boxes with actual product screenshots. If presenting live, do a live demo here instead. Have a pre-populated account ready with Week 2 or Week 3 in progress.")


# ─── SLIDE 9: TRACTION ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "What's Already Working", font_size=36, color=SLATE, bold=True)

traction_items = [
    ("100+", "weekly active users", "TimeSaveAI observation feedback tool —\nsame design philosophy, proven approach"),
    ("$5", "per teacher", "Full 7-week course cost vs.\n$50-200 for traditional PD"),
    ("7", "week course — live", "Not a prototype. Auth, onboarding, adaptive\ntutoring, artifacts, podcasts, progress tracking"),
    ("15+", "years teaching", "U.S., Australia, Finland.\nBuilt by a teacher, for teachers."),
]

for i, (num, label, desc) in enumerate(traction_items):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.6 + row * 2.6)

    tf = add_textbox(slide, left, top, Inches(2), Inches(0.8)).text_frame
    set_text(tf, num, font_size=52, color=TEAL, bold=True)
    tf2 = add_textbox(slide, Emu(left + Inches(2.2)), top, Inches(3.5), Inches(0.5)).text_frame
    set_text(tf2, label, font_size=20, color=SLATE, bold=True)
    tf3 = add_textbox(slide, Emu(left + Inches(2.2)), Emu(top + Inches(0.55)), Inches(3.5), Inches(1)).text_frame
    set_text(tf3, desc, font_size=14, color=SLATE_LIGHT)

tf = add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6)).text_frame
set_text(tf, "New Visions for Public Schools — AI literacy PD delivered to all staff, February 2026", font_size=16, color=SLATE_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Be honest about stage. 'We're pre-launch with a working product. TimeSaveAI validates the approach. The next step is a formal pilot.' Don't overclaim — authenticity is the brand.")


# ─── SLIDE 10: MARKET ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Timing + Scale", font_size=36, color=SLATE, bold=True)

# Big numbers left
market_stats = [
    ("3.7M", "K-12 teachers in the U.S."),
    ("$18B", "annual PD spend"),
    ("<1%", "is AI-specific PD"),
]

for i, (num, desc) in enumerate(market_stats):
    top = Inches(1.8 + i * 1.5)
    tf = add_textbox(slide, Inches(1), top, Inches(3), Inches(0.8)).text_frame
    set_text(tf, num, font_size=48, color=TEAL, bold=True)
    tf2 = add_textbox(slide, Inches(1), Emu(top + Inches(0.8)), Inches(4), Inches(0.5)).text_frame
    set_text(tf2, desc, font_size=18, color=SLATE_LIGHT)

# Right column - timing
add_rect(slide, Inches(6.5), Inches(1.8), Inches(6), Inches(4.5), VERY_LIGHT_TEAL)
tf = add_textbox(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.5)).text_frame
set_text(tf, "WHY NOW", font_size=16, color=TEAL_DARK, bold=True)

timing_items = [
    "DOL AI Literacy Framework (Feb 2026) creates federal compliance urgency",
    "Districts have budget for PD but no AI-specific programs to buy",
    "First-wave AI tools (MagicSchool, Khanmigo) revealed the dependency problem",
    "Category is being created — first mover with learning science wins",
    "International expansion: 70M+ teachers globally",
]
tf = add_textbox(slide, Inches(7), Inches(2.7), Inches(5.2), Inches(3.5)).text_frame
set_text(tf, timing_items[0], font_size=15, color=SLATE)
for item in timing_items[1:]:
    add_para(tf, item, font_size=15, color=SLATE, space_before=Pt(12))

add_notes(slide, "Focus on timing: 'The federal government just told every school district they need AI literacy PD. We built it before they asked.' Don't oversell TAM — the timing angle is more compelling.")


# ─── SLIDE 11: BUSINESS MODEL ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "How It Sustains", font_size=36, color=SLATE, bold=True)

phases = [
    ("PHASE 1 — NOW", "B2C Free", "Individual teachers sign up free.\nBuild user base. Collect data. Generate testimonials.", TEAL),
    ("PHASE 2 — POST-PILOT", "B2B Licensing", "District license: $25-50/teacher/year.\nAdmin dashboard, cohorts, PD credits, reporting.", AMBER),
    ("PHASE 3 — SCALE", "Platform", "Advanced modules ($10-15 each). Certification.\nCurriculum licensing. Multilingual expansion.", EMERALD),
]

for i, (phase, title, desc, color) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    add_rect(slide, left, Inches(1.6), Inches(3.6), Inches(0.06), color)
    tf = add_textbox(slide, left, Inches(1.9), Inches(3.6), Inches(0.4)).text_frame
    set_text(tf, phase, font_size=13, color=color, bold=True)
    tf2 = add_textbox(slide, left, Inches(2.3), Inches(3.6), Inches(0.5)).text_frame
    set_text(tf2, title, font_size=22, color=SLATE, bold=True)
    tf3 = add_textbox(slide, left, Inches(3.0), Inches(3.6), Inches(1.5)).text_frame
    set_text(tf3, desc, font_size=15, color=SLATE_LIGHT)

# Unit economics callout
add_rect(slide, Inches(2.5), Inches(5.0), Inches(8), Inches(1.3), VERY_LIGHT_TEAL)
tf = add_textbox(slide, Inches(3), Inches(5.15), Inches(7), Inches(0.5)).text_frame
set_text(tf, "UNIT ECONOMICS", font_size=14, color=TEAL_DARK, bold=True)
tf2 = add_textbox(slide, Inches(3), Inches(5.6), Inches(7), Inches(0.6)).text_frame
set_text(tf2, "$5 cost  /  $50 revenue  =  90% gross margin at scale", font_size=24, color=TEAL_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Key insight: the expensive part (curriculum, classifier, architecture) is already built. Marginal cost is API calls. Every new module is content, not code.")


# ─── SLIDE 12: TEAM ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Built by a Teacher, for Teachers", font_size=36, color=SLATE, bold=True)

# Photo placeholder
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(2.0), Inches(2.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
shape.line.fill.background()
tf_p = shape.text_frame
tf_p.paragraphs[0].alignment = PP_ALIGN.CENTER
set_text(tf_p, "[Photo]", font_size=18, color=SLATE_LIGHT, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(4.5), Inches(1.8), Inches(7.5), Inches(0.6)).text_frame
set_text(tf, "Asher Scott, Founder", font_size=28, color=SLATE, bold=True)

bio_items = [
    "15+ years classroom teaching — U.S., Australia, Finland",
    "Built TimeSaveAI — 100+ weekly users on observation feedback tool",
    "Delivered AI literacy PD at New Visions for Public Schools",
    "Education technologist at the intersection of pedagogy and AI",
]

tf = add_textbox(slide, Inches(4.5), Inches(2.6), Inches(7.5), Inches(3)).text_frame
set_text(tf, bio_items[0], font_size=18, color=SLATE)
for item in bio_items[1:]:
    add_para(tf, item, font_size=18, color=SLATE, space_before=Pt(14))

tf = add_textbox(slide, Inches(4.5), Inches(5.0), Inches(7.5), Inches(1.5)).text_frame
set_text(tf, '"I\'m not a tech founder who discovered education.\nI\'m a teacher who learned to build."', font_size=20, color=TEAL_DARK, italic=True)

add_notes(slide, "The credibility story: 15 years in classrooms across three countries. Built a tool with 100+ users. Delivered PD at a major NYC education org. This is teacher-first, tech-second.")


# ─── SLIDE 13: COMPETITIVE LANDSCAPE ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "We're Not Replacing AI Tools. We're the Coach.", font_size=32, color=SLATE, bold=True)

# Comparison table
headers = ["", "Skippy", "Khanmigo", "MagicSchool", "ChatGPT", "Traditional PD"]
col_widths = [Inches(2.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8)]
col_lefts = [Inches(0.6)]
for w in col_widths[:-1]:
    col_lefts.append(Emu(col_lefts[-1] + w))

# Header row
top = Inches(1.5)
add_rect(slide, Inches(0.6), top, Inches(11.8), Inches(0.5), TEAL)
for j, h in enumerate(headers):
    tf = add_textbox(slide, col_lefts[j], top, col_widths[j], Inches(0.5)).text_frame
    set_text(tf, h, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rows = [
    ("Teaches capability", "Yes", "No", "No", "No", "Sometimes"),
    ("Personalized", "6 layers", "No", "Partial", "No", "No"),
    ("Structured progression", "7 weeks", "No", "No", "No", "1 session"),
    ("Reusable artifacts", "7", "No", "Yes*", "Ad hoc", "Sometimes"),
    ("Builds independence", "Yes", "No", "No", "N/A", "Varies"),
    ("Cost/teacher", "~$5", "$44/yr", "$99/yr", "Free-$20/mo", "$50-200"),
]

for i, row in enumerate(rows):
    row_top = Emu(Inches(2.05) + i * Inches(0.65))
    bg = VERY_LIGHT_TEAL if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.6), row_top, Inches(11.8), Inches(0.6), bg)
    # Highlight Skippy column
    if i % 2 == 1:
        add_rect(slide, col_lefts[1], row_top, col_widths[1], Inches(0.6), VERY_LIGHT_TEAL)

    for j, val in enumerate(row):
        color = TEAL_DARK if j == 1 and val.startswith(("Yes", "6", "7", "~")) else SLATE
        b = j == 1
        tf = add_textbox(slide, col_lefts[j], row_top, col_widths[j], Inches(0.6)).text_frame
        set_text(tf, val, font_size=13, color=color, bold=b, alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

tf = add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.6)).text_frame
set_text(tf, "* MagicSchool artifacts are tool-generated, not teacher-created. No skill transfer.", font_size=13, color=SLATE_LIGHT, italic=True)

add_notes(slide, "Don't trash competitors. Frame it as: 'These tools serve different purposes. Khanmigo tutors students. MagicSchool generates content. We build teacher capability. They're complements, not substitutes.'")


# ─── SLIDE 14: ROADMAP ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "Where We're Going", font_size=36, color=SLATE, bold=True)

# Timeline
milestones = [
    ("NOW", "Launch & Polish", "Streaming responses\nArtifact export\nError resilience\n10-school pilot", TEAL),
    ("3 MONTHS", "Admin & Scale", "Admin dashboard\nPD credit certification\nCohort management\nDistrict outreach", AMBER),
    ("6 MONTHS", "Expand", "Advanced modules\nVoice mode\nMultilingual (Spanish)\nResearch publication", EMERALD),
    ("12 MONTHS", "Platform", "Curriculum licensing\nDistrict API integrations\nNational expansion\nEfficacy data", RGBColor(0x63, 0x66, 0xF1)),
]

# Timeline line
add_rect(slide, Inches(1.5), Inches(2.8), Inches(10.5), Inches(0.04), RGBColor(0xE2, 0xE8, 0xF0))

for i, (time, title, items, color) in enumerate(milestones):
    left = Inches(1.2 + i * 2.9)
    # Dot on timeline
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(left + Inches(1.1)), Inches(2.65), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    tf = add_textbox(slide, left, Inches(1.8), Inches(2.6), Inches(0.4)).text_frame
    set_text(tf, time, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    tf2 = add_textbox(slide, left, Inches(2.15), Inches(2.6), Inches(0.5)).text_frame
    set_text(tf2, title, font_size=17, color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)
    tf3 = add_textbox(slide, left, Inches(3.3), Inches(2.6), Inches(2.5)).text_frame
    set_text(tf3, items, font_size=14, color=SLATE_LIGHT, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6)).text_frame
set_text(tf, "The architecture is built. What's left is polish, scale, and data.", font_size=18, color=SLATE, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "The point: the hard engineering work is done. The roadmap is about refining the experience, collecting evidence, and scaling distribution.")


# ─── SLIDE 15: IMPACT VISION ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, SLATE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.8)).text_frame
set_text(tf, "The World We're Building", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

vision_items = [
    "Every teacher in America has access to personalized AI coaching — for $5",
    "Teacher AI literacy becomes measurable, not anecdotal",
    "The profession develops its relationship with AI on its own terms",
]

tf = add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(2.5)).text_frame
set_text(tf, vision_items[0], font_size=22, color=TEAL_LIGHT)
for item in vision_items[1:]:
    add_para(tf, item, font_size=22, color=TEAL_LIGHT, space_before=Pt(18))

# Quote
tf = add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1.5)).text_frame
set_text(tf, '"I think the question isn\'t whether you can use AI.\nIt\'s whether you can use AI and still be\nrecognizably yourself at the end of it."', font_size=24, color=WHITE, italic=True, alignment=PP_ALIGN.CENTER)

tf2 = add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5)).text_frame
set_text(tf2, "— Asher Scott", font_size=16, color=TEAL, alignment=PP_ALIGN.CENTER)

add_notes(slide, "This is the emotional close. Slow down. Make eye contact. Let the quote land. This is from Asher's own PD presentation at New Visions.")


# ─── SLIDE 16: THE ASK ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), TEAL)

tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8)).text_frame
set_text(tf, "What We Need", font_size=36, color=SLATE, bold=True)

asks = [
    ("Funding", "$500K seed (or $150-250K grant)", TEAL),
    ("Pilots", "10 schools, 200 teachers, full data collection", AMBER),
    ("Partnerships", "Districts ready to lead on AI literacy PD", EMERALD),
    ("Introductions", "Education foundations, accelerators, district CTOs", RGBColor(0x63, 0x66, 0xF1)),
]

for i, (title, desc, color) in enumerate(asks):
    top = Inches(1.6 + i * 1.15)
    add_rect(slide, Inches(0.8), top, Inches(0.4), Inches(0.7), color)
    tf = add_textbox(slide, Inches(1.5), top, Inches(4), Inches(0.4)).text_frame
    set_text(tf, title, font_size=22, color=SLATE, bold=True)
    tf2 = add_textbox(slide, Inches(1.5), Emu(top + Inches(0.4)), Inches(8), Inches(0.5)).text_frame
    set_text(tf2, desc, font_size=17, color=SLATE_LIGHT)

# Contact
add_rect(slide, Inches(3), Inches(5.8), Inches(7), Inches(1.2), VERY_LIGHT_TEAL)
tf = add_textbox(slide, Inches(3.5), Inches(5.9), Inches(6), Inches(0.5)).text_frame
set_text(tf, "Asher Scott", font_size=22, color=SLATE, bold=True, alignment=PP_ALIGN.CENTER)
tf2 = add_textbox(slide, Inches(3.5), Inches(6.4), Inches(6), Inches(0.5)).text_frame
set_text(tf2, "asher.rhys.scott@gmail.com", font_size=17, color=TEAL_DARK, alignment=PP_ALIGN.CENTER)

tf = add_textbox(slide, Inches(2), Inches(6.9), Inches(9), Inches(0.5)).text_frame
set_text(tf, "AI for Teachers: Capacity, not dependency.", font_size=16, color=SLATE_LIGHT, italic=True, alignment=PP_ALIGN.CENTER)

add_notes(slide, "Be specific about what you're asking for. End with: 'I'd love to show you the product. Can we schedule 15 minutes this week?'")


# ─── SAVE ──────────────────────────────────────────────────────
output_path = "/Users/carladelporto/ai-for-teachers/docs/pitch/decks/AI-for-Teachers-Master-Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
